from pathlib import Path

TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".log", ".json", ".xml", ".html", ".htm"}
PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx"}
PPTX_EXTENSIONS = {".pptx"}
XLSX_EXTENSIONS = {".xlsx", ".xls"}

MAX_CHARS = 500


def extract_text(path: Path) -> str | None:
    ext = path.suffix.lower()
    try:
        if ext in TEXT_EXTENSIONS:
            return path.read_text(encoding="utf-8", errors="ignore")[:MAX_CHARS]
        if ext in PDF_EXTENSIONS:
            return _read_pdf(path)
        if ext in DOCX_EXTENSIONS:
            return _read_docx(path)
        if ext in PPTX_EXTENSIONS:
            return _read_pptx(path)
        if ext in XLSX_EXTENSIONS:
            return _read_xlsx(path)
    except Exception:
        return None
    return None


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
        if len(text) >= MAX_CHARS:
            break
    return text[:MAX_CHARS]


def _read_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    text = "\n".join(p.text for p in doc.paragraphs)
    return text[:MAX_CHARS]


def _read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)[:MAX_CHARS]


def _read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            parts.append(" ".join(str(c) for c in row if c is not None))
            if len("\n".join(parts)) >= MAX_CHARS:
                break
        break
    return "\n".join(parts)[:MAX_CHARS]
